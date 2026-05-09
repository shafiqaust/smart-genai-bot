from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color palette ──
BG_DARK = RGBColor(0x1B, 0x1B, 0x2F)
BG_CARD = RGBColor(0x27, 0x27, 0x45)
BG_CARD_LIGHT = RGBColor(0x32, 0x32, 0x55)
ACCENT = RGBColor(0x6C, 0x63, 0xFF)
ACCENT2 = RGBColor(0x00, 0xD2, 0xFF)
ACCENT3 = RGBColor(0xFF, 0x6B, 0x6B)
ACCENT4 = RGBColor(0x4E, 0xCB, 0x71)
ORANGE = RGBColor(0xFF, 0xA5, 0x00)
YELLOW = RGBColor(0xFF, 0xD9, 0x3D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
MID_GRAY = RGBColor(0x88, 0x88, 0xAA)
DARK_ACCENT = RGBColor(0x3D, 0x38, 0x99)


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, corner_radius=None):
    if corner_radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.adjustments[0] = corner_radius
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_circle(slide, left, top, size, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_arrow_right(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_arrow_down(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_chevron(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_pentagon(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.PENTAGON, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def icon_circle_text(slide, left, top, size, bg_color, text, font_size=28):
    circle = add_circle(slide, left, top, size, bg_color)
    tf = circle.text_frame
    tf.paragraphs[0].text = text
    tf.paragraphs[0].font.size = Pt(font_size)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return circle


def add_flow_step(slide, x, y, w, h, num, title, desc, color, num_size=Inches(0.7)):
    card = add_shape(slide, x, y, w, h, BG_CARD, corner_radius=0.05)
    add_shape(slide, x, y, w, Inches(0.06), color)
    icon_circle_text(slide, x + w/2 - num_size/2, y + Inches(0.2), num_size, color, num, font_size=22)
    add_text_box(slide, x + Inches(0.1), y + num_size + Inches(0.25), w - Inches(0.2), Inches(0.5),
                 title, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), y + num_size + Inches(0.7), w - Inches(0.2), h - num_size - Inches(0.9),
                 desc, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    return card


# ════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT)

# Decorative circles
add_circle(slide, Inches(-0.3), Inches(-0.3), Inches(1.5), DARK_ACCENT)
add_circle(slide, Inches(12.3), Inches(6.3), Inches(1.5), DARK_ACCENT)

# Icon row
icons = [("🎓", ACCENT), ("🤖", ACCENT2), ("📊", ACCENT3), ("🔍", ACCENT4)]
for i, (emoji, color) in enumerate(icons):
    icon_circle_text(slide, Inches(4.2) + i * Inches(1.4), Inches(1.4), Inches(1.0), color, emoji, font_size=32)

add_text_box(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(1.2),
             "Smart Academic Assistant", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(5.2), Inches(4.0), Inches(2.9), Inches(0.04), ACCENT)

add_text_box(slide, Inches(1.5), Inches(4.3), Inches(10), Inches(0.8),
             "A Multi-Tool RAG System for Student Advising", font_size=24, color=ACCENT2, alignment=PP_ALIGN.CENTER)

# Tech badges
techs = [("LangGraph", ACCENT), ("GPT-4.1", ACCENT2), ("RAG", ACCENT3), ("FastAPI", ACCENT4)]
for i, (label, color) in enumerate(techs):
    x = Inches(3.5) + i * Inches(1.8)
    badge = add_shape(slide, x, Inches(5.5), Inches(1.5), Inches(0.5), color, corner_radius=0.1)
    tf = badge.text_frame
    tf.paragraphs[0].text = label
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


# ════════════════════════════════════════════
# SLIDE 2 — The Problem
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "The Problem", font_size=38, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.05), ACCENT)

problems = [
    ("💬", "Understanding\nLanguage", "Student questions are messy\nand require reasoning", ACCENT),
    ("📎", "Processing\nUploads", "PDFs, URLs, text — each\nneeds different handling", ACCENT2),
    ("🎯", "Accurate\nKnowledge", "LLMs hallucinate — answers\nmust be grounded in data", ACCENT3),
    ("💡", "Clear\nExplanations", "Correct AND clearly\nexplained answers", ACCENT4),
]

for i, (emoji, title, desc, color) in enumerate(problems):
    x = Inches(0.6) + i * Inches(3.15)
    card = add_shape(slide, x, Inches(1.8), Inches(2.85), Inches(4.8), BG_CARD, corner_radius=0.06)

    icon_circle_text(slide, x + Inches(0.85), Inches(2.15), Inches(1.15), color, emoji, font_size=36)

    add_text_box(slide, x + Inches(0.15), Inches(3.5), Inches(2.55), Inches(0.9),
                 title, font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_shape(slide, x + Inches(0.8), Inches(4.4), Inches(1.25), Inches(0.03), color)
    add_text_box(slide, x + Inches(0.15), Inches(4.6), Inches(2.55), Inches(1.5),
                 desc, font_size=15, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════
# SLIDE 3 — Solution Overview (Visual Pipeline)
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "LangGraph-Powered Architecture", font_size=38, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.05), ACCENT)

# LangGraph orchestration banner
lg_banner = add_shape(slide, Inches(0.4), Inches(1.4), Inches(12.5), Inches(0.65), ACCENT, corner_radius=0.08)
tf = lg_banner.text_frame
tf.paragraphs[0].text = "🔀  LangGraph StateGraph orchestrates the ENTIRE system — every stage, every decision"
tf.paragraphs[0].font.size = Pt(16)
tf.paragraphs[0].font.color.rgb = WHITE
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

# User icon on left
icon_circle_text(slide, Inches(0.5), Inches(3.2), Inches(1.3), MID_GRAY, "👤", font_size=40)
add_text_box(slide, Inches(0.2), Inches(4.6), Inches(1.7), Inches(0.5),
             "Student", font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Arrow from user
add_arrow_right(slide, Inches(1.9), Inches(3.6), Inches(0.8), Inches(0.4), MID_GRAY)

# LangGraph frame around all stages
lg_frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(2.05), Inches(10.6), Inches(4.35))
lg_frame.adjustments[0] = 0.03
lg_frame.fill.solid()
lg_frame.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x38)
lg_frame.line.color.rgb = ACCENT
lg_frame.line.width = Pt(2)

# Stage 1
s1 = add_shape(slide, Inches(2.8), Inches(2.2), Inches(2.8), Inches(3.2), BG_CARD, corner_radius=0.06)
add_shape(slide, Inches(2.8), Inches(2.2), Inches(2.8), Inches(0.07), ACCENT)
icon_circle_text(slide, Inches(3.6), Inches(2.45), Inches(1.2), ACCENT, "📥", font_size=36)
add_text_box(slide, Inches(2.9), Inches(3.8), Inches(2.6), Inches(0.5),
             "Input Processing", font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(2.9), Inches(4.3), Inches(2.6), Inches(0.8),
             "Parse uploads\nExtract structured data", font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Arrow
add_arrow_right(slide, Inches(5.7), Inches(3.6), Inches(0.8), Inches(0.4), ACCENT)

# Stage 2
s2 = add_shape(slide, Inches(6.6), Inches(2.2), Inches(2.8), Inches(3.2), BG_CARD, corner_radius=0.06)
add_shape(slide, Inches(6.6), Inches(2.2), Inches(2.8), Inches(0.07), ACCENT2)
icon_circle_text(slide, Inches(7.4), Inches(2.45), Inches(1.2), ACCENT2, "⚙️", font_size=36)
add_text_box(slide, Inches(6.7), Inches(3.8), Inches(2.6), Inches(0.5),
             "Data Processing", font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(6.7), Inches(4.3), Inches(2.6), Inches(0.8),
             "Classify, retrieve\nCalculate", font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Arrow
add_arrow_right(slide, Inches(9.5), Inches(3.6), Inches(0.8), Inches(0.4), ACCENT2)

# Stage 3
s3 = add_shape(slide, Inches(10.4), Inches(2.2), Inches(2.5), Inches(3.2), BG_CARD, corner_radius=0.06)
add_shape(slide, Inches(10.4), Inches(2.2), Inches(2.5), Inches(0.07), ACCENT4)
icon_circle_text(slide, Inches(11.05), Inches(2.45), Inches(1.2), ACCENT4, "✨", font_size=36)
add_text_box(slide, Inches(10.5), Inches(3.8), Inches(2.3), Inches(0.5),
             "Answer Gen", font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(10.5), Inches(4.3), Inches(2.3), Inches(0.8),
             "Grounded response\nvia GPT-4.1", font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# File labels
add_text_box(slide, Inches(2.9), Inches(5.5), Inches(2.6), Inches(0.4),
             "server.py", font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(6.7), Inches(5.5), Inches(2.6), Inches(0.4),
             "rag_bot.py", font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(10.5), Inches(5.5), Inches(2.3), Inches(0.4),
             "rag_bot.py", font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# LangGraph node labels below file labels
add_text_box(slide, Inches(2.9), Inches(5.85), Inches(2.6), Inches(0.4),
             "LangGraph orchestrated", font_size=11, color=ACCENT, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(6.7), Inches(5.85), Inches(2.6), Inches(0.4),
             "LangGraph nodes", font_size=11, color=ACCENT, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(10.5), Inches(5.85), Inches(2.3), Inches(0.4),
             "LangGraph node", font_size=11, color=ACCENT, alignment=PP_ALIGN.CENTER)

# LangGraph frame label
add_text_box(slide, Inches(2.5), Inches(6.45), Inches(10.6), Inches(0.4),
             "▲  LangGraph StateGraph  ▲", font_size=13, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════
# SLIDE — LangGraph: The Brain of the System
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "LangGraph — The Brain", font_size=38, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.05), ACCENT)

add_text_box(slide, Inches(0.8), Inches(1.35), Inches(11), Inches(0.5),
             "Every decision, every data flow, every tool call is orchestrated by a LangGraph StateGraph",
             font_size=16, color=ACCENT2)

# Graph visualization: 5 processing nodes
lg_nodes = [
    ("load_context", "📂", "Load catalog +\ntranscript JSON", ACCENT),
    ("classify_intent", "🏷️", "Route question to\ncalc / reason / RAG", ACCENT2),
    ("retrieve_docs", "🔍", "Semantic search\nover vector KB", ACCENT3),
    ("calculate_tools", "🧮", "CGPA, credits,\ngrad status", ACCENT4),
    ("generate_answer", "✨", "Assemble all context\n→ GPT-4.1 response", ORANGE),
]

for i, (name, emoji, desc, color) in enumerate(lg_nodes):
    x = Inches(0.3) + i * Inches(2.6)
    y = Inches(2.1)
    card = add_shape(slide, x, y, Inches(2.3), Inches(2.7), BG_CARD, corner_radius=0.06)
    add_shape(slide, x, y, Inches(2.3), Inches(0.06), color)
    icon_circle_text(slide, x + Inches(0.55), y + Inches(0.2), Inches(1.2), color, emoji, font_size=32)
    add_text_box(slide, x + Inches(0.1), y + Inches(1.5), Inches(2.1), Inches(0.45),
                 name, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), y + Inches(1.95), Inches(2.1), Inches(0.6),
                 desc, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    if i < 4:
        add_arrow_right(slide, x + Inches(2.35), y + Inches(0.65), Inches(0.3), Inches(0.2), color)

# Shared State section
state_card = add_shape(slide, Inches(0.8), Inches(5.1), Inches(11.5), Inches(2.2), BG_CARD, corner_radius=0.06)
add_shape(slide, Inches(0.8), Inches(5.1), Inches(11.5), Inches(0.06), ACCENT)
add_text_box(slide, Inches(1.2), Inches(5.2), Inches(5), Inches(0.5),
             "🔀  LangGraph Shared State (passed between every node)", font_size=18, color=ACCENT, bold=True)

state_items = [
    ("💬", "question — the student's original message"),
    ("📋", "catalog_rules — program reqs from URL"),
    ("📄", "transcript_record — student data from upload"),
    ("🏷️", "intent_route — calc / reasoning / retrieval / hybrid"),
    ("🔍", "retrieved_chunks + citations — RAG results"),
    ("🧮", "calculation_results — CGPA, credits, grad status"),
]
for i, (icon, text) in enumerate(state_items):
    col = i % 2
    row = i // 2
    x = Inches(1.2) + col * Inches(5.8)
    y = Inches(5.8) + row * Inches(0.48)
    add_text_box(slide, x, y, Inches(0.4), Inches(0.4), icon, font_size=13)
    add_text_box(slide, x + Inches(0.4), y + Inches(0.02), Inches(5.0), Inches(0.4),
                 text, font_size=12, color=LIGHT_GRAY)


# ════════════════════════════════════════════
# SLIDE 4 — Stage 1: Section Title
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT)

# Big icon
icon_circle_text(slide, Inches(5.65), Inches(1.2), Inches(2.0), ACCENT, "📥", font_size=64)

add_text_box(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(0.5),
             "STAGE I", font_size=22, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(4.0), Inches(10), Inches(1.0),
             "Input Processing", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(5.2), Inches(5.0), Inches(2.9), Inches(0.04), ACCENT)
add_text_box(slide, Inches(1.5), Inches(5.3), Inches(10), Inches(0.6),
             "Format-based routing — what you send determines how it's processed",
             font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Five input type icons
input_icons = [("🔗", "URL", ACCENT), ("📄", "PDF", ACCENT2), ("🖼️", "Image", ACCENT3), ("📋", "Text Paste", ORANGE), ("💬", "Chat", ACCENT4)]
for i, (emoji, label, color) in enumerate(input_icons):
    x = Inches(1.2) + i * Inches(2.2)
    icon_circle_text(slide, x + Inches(0.3), Inches(6.0), Inches(0.8), color, emoji, font_size=24)
    add_text_box(slide, x - Inches(0.2), Inches(6.85), Inches(1.8), Inches(0.4),
                 label, font_size=13, color=color, bold=True, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════
# SLIDE 4a — URL Input (Visual Flow)
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8),
             "Input: URL", font_size=36, color=WHITE, bold=True)
tag = add_shape(slide, Inches(5.5), Inches(0.45), Inches(2.5), Inches(0.55), ACCENT, corner_radius=0.1)
tf = tag.text_frame
tf.paragraphs[0].text = "→ catalog_rules.json"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.color.rgb = WHITE
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

add_shape(slide, Inches(0.8), Inches(1.15), Inches(2), Inches(0.05), ACCENT)

# Visual flow: URL icon → Scrape → Clean → GPT-4.1 → JSON → Disk
flow_items = [
    ("🔗", "URL\nDetected", "Starts with\nhttp(s)://"),
    ("🌐", "Scrape\nWebpage", "HTTP GET with\nbrowser headers"),
    ("🧹", "Clean\nHTML", "Strip tags,\nextract text"),
    ("🤖", "GPT-4.1\nExtract", "Structured output\nw/ JSON schema"),
    ("💾", "Save\nto Disk", "catalog_rules.json\nin extracted_data/"),
]

for i, (emoji, title, desc) in enumerate(flow_items):
    x = Inches(0.4) + i * Inches(2.55)
    card = add_shape(slide, x, Inches(1.7), Inches(2.2), Inches(3.0), BG_CARD, corner_radius=0.06)
    icon_circle_text(slide, x + Inches(0.6), Inches(1.9), Inches(1.0), ACCENT if i != 3 else ACCENT2, emoji, font_size=30)
    add_text_box(slide, x + Inches(0.1), Inches(3.05), Inches(2.0), Inches(0.6),
                 title, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), Inches(3.65), Inches(2.0), Inches(0.8),
                 desc, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    if i < 4:
        add_arrow_right(slide, x + Inches(2.25), Inches(2.85), Inches(0.35), Inches(0.25), MID_GRAY)

# ════════════════════════════════════════════
# SLIDE — URL Input Example
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "Input: URL — Example", font_size=36, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.15), Inches(2), Inches(0.05), ACCENT)

ex_card = add_shape(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(5.0), BG_CARD, corner_radius=0.05)
add_shape(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.06), ACCENT)
add_text_box(slide, Inches(1.2), Inches(1.9), Inches(3), Inches(0.5),
             "💡  Real Example", font_size=22, color=ACCENT, bold=True)

# Input side
in_card = add_shape(slide, Inches(1.0), Inches(2.6), Inches(4.8), Inches(2.5), BG_CARD_LIGHT, corner_radius=0.05)
add_text_box(slide, Inches(1.2), Inches(2.7), Inches(1.0), Inches(0.35),
             "INPUT", font_size=14, color=ACCENT2, bold=True)
add_text_box(slide, Inches(1.2), Inches(3.3), Inches(4.4), Inches(1.2),
             "https://catalogs.nmsu.edu/nmsu/\n  arts-sciences/computer-science/",
             font_size=16, color=LIGHT_GRAY)

# Arrow
add_arrow_right(slide, Inches(6.0), Inches(3.5), Inches(0.6), Inches(0.3), ACCENT)

# Output side
out_card = add_shape(slide, Inches(6.8), Inches(2.6), Inches(5.3), Inches(3.8), BG_CARD_LIGHT, corner_radius=0.05)
add_text_box(slide, Inches(7.0), Inches(2.7), Inches(1.2), Inches(0.35),
             "OUTPUT", font_size=14, color=ACCENT4, bold=True)
add_text_box(slide, Inches(7.0), Inches(3.2), Inches(4.9), Inches(3.0),
             '{ "program_name": "Computer Science",\n  "required_credits": 120,\n  "minimum_cgpa_for_graduation": 2.0,\n  "tracks": [\n    {"track_name": "AI",\n     "required_credits": 15, "notes": "..."}],\n  "required_courses": [\n    {"course_name": "CS 172",\n     "credits": 4, "category": "Core"}],\n  "elective_credit_requirement": 12,\n  "exceptions_and_notes":\n    ["B.S. is ABET accredited"],\n  "source_url": "https://catalogs..." }',
             font_size=12, color=LIGHT_GRAY)


# ════════════════════════════════════════════
# SLIDE 4b — PDF Input
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8),
             "Input: PDF", font_size=36, color=WHITE, bold=True)
tag = add_shape(slide, Inches(5.0), Inches(0.45), Inches(3.2), Inches(0.55), ACCENT2, corner_radius=0.1)
tf = tag.text_frame
tf.paragraphs[0].text = "→ transcript_record.json"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.color.rgb = WHITE
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
add_shape(slide, Inches(0.8), Inches(1.15), Inches(2), Inches(0.05), ACCENT2)

# Visual flow: Upload → Save temp → PDFReader → GPT-4.1 → Save JSON
pdf_flow_items = [
    ("📎", "File\nUpload", "User uploads\na PDF file"),
    ("💾", "Save\nTemp", "Write to temp\ndirectory"),
    ("📄", "PDFReader\nExtract", "Extract raw text\nfrom PDF pages"),
    ("🤖", "GPT-4.1\nParse", "Structured output\nw/ JSON schema"),
    ("💾", "Save\nto Disk", "transcript_record.json\nin extracted_data/"),
]

for i, (emoji, title, desc) in enumerate(pdf_flow_items):
    x = Inches(0.4) + i * Inches(2.55)
    card = add_shape(slide, x, Inches(1.7), Inches(2.2), Inches(3.0), BG_CARD, corner_radius=0.06)
    icon_circle_text(slide, x + Inches(0.6), Inches(1.9), Inches(1.0), ACCENT2 if i != 3 else ACCENT, emoji, font_size=30)
    add_text_box(slide, x + Inches(0.1), Inches(3.05), Inches(2.0), Inches(0.6),
                 title, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), Inches(3.65), Inches(2.0), Inches(0.8),
                 desc, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    if i < 4:
        add_arrow_right(slide, x + Inches(2.25), Inches(2.85), Inches(0.35), Inches(0.25), MID_GRAY)

# ════════════════════════════════════════════
# SLIDE — PDF Example
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "Input: PDF — Example", font_size=36, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.15), Inches(2), Inches(0.05), ACCENT2)

ex_card = add_shape(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(5.0), BG_CARD, corner_radius=0.05)
add_shape(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.06), ACCENT4)
add_text_box(slide, Inches(1.2), Inches(1.9), Inches(3), Inches(0.5),
             "💡  Real Example Output", font_size=22, color=ACCENT4, bold=True)

# Input side
in_card = add_shape(slide, Inches(1.0), Inches(2.6), Inches(4.8), Inches(2.5), BG_CARD_LIGHT, corner_radius=0.05)
add_text_box(slide, Inches(1.2), Inches(2.7), Inches(1.0), Inches(0.35),
             "INPUT", font_size=14, color=ACCENT2, bold=True)
icon_circle_text(slide, Inches(1.8), Inches(3.3), Inches(1.0), ACCENT2, "📄", font_size=30)
add_text_box(slide, Inches(3.0), Inches(3.45), Inches(2.6), Inches(0.5),
             "transcript.pdf\n(uploaded)", font_size=16, color=LIGHT_GRAY)

# Arrow
add_arrow_right(slide, Inches(6.0), Inches(3.5), Inches(0.6), Inches(0.3), ACCENT4)

# Output side
out_card = add_shape(slide, Inches(6.8), Inches(2.6), Inches(5.3), Inches(3.8), BG_CARD_LIGHT, corner_radius=0.05)
add_text_box(slide, Inches(7.0), Inches(2.7), Inches(1.2), Inches(0.35),
             "OUTPUT", font_size=14, color=ACCENT4, bold=True)
add_text_box(slide, Inches(7.0), Inches(3.2), Inches(4.9), Inches(2.8),
             '{ "student_name": "Albert Smith",\n  "program": "Computer Science",\n  "completed_courses": [\n    {"course": "Chemistry",\n     "credits": 1, "grade": "A"},\n    {"course": "Calculus I",\n     "credits": 4, "grade": "B+"},\n    ... ] }',
             font_size=13, color=LIGHT_GRAY)

# ════════════════════════════════════════════
# SLIDE 4b2 — Image Input
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8),
             "Input: Image", font_size=36, color=WHITE, bold=True)
tag = add_shape(slide, Inches(5.3), Inches(0.45), Inches(3.2), Inches(0.55), ACCENT3, corner_radius=0.1)
tf = tag.text_frame
tf.paragraphs[0].text = "→ transcript_record.json"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.color.rgb = WHITE
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
add_shape(slide, Inches(0.8), Inches(1.15), Inches(2), Inches(0.05), ACCENT3)

# Visual flow: Upload → Base64 → Data URL → GPT-4.1 Vision → Save JSON
img_flow_items = [
    ("📎", "File\nUpload", "User uploads\nan image file"),
    ("🔢", "Base64\nEncode", "Convert binary\nto base64 string"),
    ("🔗", "Build\nData URL", "Create data URI\nfor API payload"),
    ("🤖", "GPT-4.1\nVision", "Multimodal parse\nw/ JSON schema"),
    ("💾", "Save\nto Disk", "transcript_record.json\nin extracted_data/"),
]

for i, (emoji, title, desc) in enumerate(img_flow_items):
    x = Inches(0.4) + i * Inches(2.55)
    card = add_shape(slide, x, Inches(1.7), Inches(2.2), Inches(3.0), BG_CARD, corner_radius=0.06)
    icon_circle_text(slide, x + Inches(0.6), Inches(1.9), Inches(1.0), ACCENT3 if i != 3 else ACCENT, emoji, font_size=30)
    add_text_box(slide, x + Inches(0.1), Inches(3.05), Inches(2.0), Inches(0.6),
                 title, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), Inches(3.65), Inches(2.0), Inches(0.8),
                 desc, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    if i < 4:
        add_arrow_right(slide, x + Inches(2.25), Inches(2.85), Inches(0.35), Inches(0.25), MID_GRAY)

# ════════════════════════════════════════════
# SLIDE — Image Example
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "Input: Image — Example", font_size=36, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.15), Inches(2), Inches(0.05), ACCENT3)

ex_card = add_shape(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(5.0), BG_CARD, corner_radius=0.05)
add_shape(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.06), ACCENT4)
add_text_box(slide, Inches(1.2), Inches(1.9), Inches(3), Inches(0.5),
             "💡  Real Example Output", font_size=22, color=ACCENT4, bold=True)

# Input side
in_card = add_shape(slide, Inches(1.0), Inches(2.6), Inches(4.8), Inches(2.5), BG_CARD_LIGHT, corner_radius=0.05)
add_text_box(slide, Inches(1.2), Inches(2.7), Inches(1.0), Inches(0.35),
             "INPUT", font_size=14, color=ACCENT2, bold=True)
icon_circle_text(slide, Inches(1.8), Inches(3.3), Inches(1.0), ACCENT3, "🖼️", font_size=30)
add_text_box(slide, Inches(3.0), Inches(3.45), Inches(2.6), Inches(0.5),
             "transcript_image.png\n(uploaded)", font_size=16, color=LIGHT_GRAY)

# Arrow
add_arrow_right(slide, Inches(6.0), Inches(3.5), Inches(0.6), Inches(0.3), ACCENT4)

# Output side
out_card = add_shape(slide, Inches(6.8), Inches(2.6), Inches(5.3), Inches(3.8), BG_CARD_LIGHT, corner_radius=0.05)
add_text_box(slide, Inches(7.0), Inches(2.7), Inches(1.2), Inches(0.35),
             "OUTPUT", font_size=14, color=ACCENT4, bold=True)
add_text_box(slide, Inches(7.0), Inches(3.2), Inches(4.9), Inches(2.8),
             '{ "student_name": "Albert Smith",\n  "program": "Computer Science",\n  "completed_courses": [\n    {"course": "Chemistry",\n     "credits": 1, "grade": "A"},\n    {"course": "Calculus I",\n     "credits": 4, "grade": "B+"},\n    ... ] }',
             font_size=13, color=LIGHT_GRAY)



# ════════════════════════════════════════════
# SLIDE 4c — Long Text Paste
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8),
             "Input: Long Text Paste", font_size=36, color=WHITE, bold=True)
tag = add_shape(slide, Inches(6.5), Inches(0.45), Inches(3.2), Inches(0.55), ACCENT3, corner_radius=0.1)
tf = tag.text_frame
tf.paragraphs[0].text = "→ transcript_record.json"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.color.rgb = WHITE
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
add_shape(slide, Inches(0.8), Inches(1.15), Inches(2), Inches(0.05), ACCENT3)

# Detection heuristic visual
detect_card = add_shape(slide, Inches(0.8), Inches(1.7), Inches(5.2), Inches(3.5), BG_CARD, corner_radius=0.06)
add_text_box(slide, Inches(1.2), Inches(1.9), Inches(4.4), Inches(0.5),
             "🔍  Detection Heuristic", font_size=22, color=ACCENT3, bold=True)

# Visual checklist
checks = [
    ("✅", "len(message) > 300 characters", ACCENT4),
    ("✅", 'Contains "\\n" (newline)', ACCENT4),
    ("➡️", "Treated as pasted transcript", ACCENT3),
]
for i, (icon, text, color) in enumerate(checks):
    y = Inches(2.6) + i * Inches(0.8)
    icon_circle_text(slide, Inches(1.3), y, Inches(0.55), color, icon, font_size=18)
    add_text_box(slide, Inches(2.1), y + Inches(0.08), Inches(3.5), Inches(0.5),
                 text, font_size=16, color=LIGHT_GRAY)

# Flow diagram on right
flow_card = add_shape(slide, Inches(6.5), Inches(1.7), Inches(6.0), Inches(3.5), BG_CARD, corner_radius=0.06)
add_text_box(slide, Inches(6.9), Inches(1.9), Inches(5.2), Inches(0.5),
             "⚡  Processing Flow", font_size=22, color=ACCENT, bold=True)

flow_steps = [
    ("📋", "Raw Text"),
    ("🤖", "GPT-4.1"),
    ("📊", "JSON Schema"),
    ("💾", "Save"),
]
for i, (emoji, label) in enumerate(flow_steps):
    x = Inches(6.8) + i * Inches(1.4)
    icon_circle_text(slide, x, Inches(2.8), Inches(0.9), ACCENT if i != 1 else ACCENT2, emoji, font_size=26)
    add_text_box(slide, x - Inches(0.15), Inches(3.8), Inches(1.2), Inches(0.5),
                 label, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    if i < 3:
        add_arrow_right(slide, x + Inches(0.95), Inches(3.05), Inches(0.35), Inches(0.2), MID_GRAY)

# ════════════════════════════════════════════
# SLIDE — Long Text Paste Example
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "Input: Long Text Paste — Example", font_size=36, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.15), Inches(2), Inches(0.05), ACCENT3)

ex_card = add_shape(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(3.8), BG_CARD, corner_radius=0.05)
add_shape(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.06), ACCENT4)
add_text_box(slide, Inches(1.2), Inches(1.9), Inches(3), Inches(0.5),
             "💡  Example", font_size=22, color=ACCENT4, bold=True)

in_card = add_shape(slide, Inches(1.0), Inches(2.6), Inches(5.0), Inches(2.5), BG_CARD_LIGHT, corner_radius=0.05)
add_text_box(slide, Inches(1.2), Inches(2.7), Inches(1.0), Inches(0.35),
             "INPUT", font_size=14, color=ACCENT2, bold=True)
add_text_box(slide, Inches(1.2), Inches(3.2), Inches(4.6), Inches(1.5),
             '"Honors Eng 9: World Lit  A  1\\nAlg II/Trig  A  2\\nHonors Precalculus  A  1\\n..."  (450 chars)',
             font_size=15, color=LIGHT_GRAY)

add_arrow_right(slide, Inches(6.2), Inches(3.5), Inches(0.6), Inches(0.3), ACCENT4)

out_card = add_shape(slide, Inches(7.0), Inches(2.6), Inches(5.1), Inches(2.5), BG_CARD_LIGHT, corner_radius=0.05)
add_text_box(slide, Inches(7.2), Inches(2.7), Inches(1.2), Inches(0.35),
             "OUTPUT", font_size=14, color=ACCENT4, bold=True)
add_text_box(slide, Inches(7.2), Inches(3.2), Inches(4.7), Inches(1.5),
             '{ "student_name": "Albert Smith",\n  "program": null,\n  "completed_courses": [\n    {"course": "Honors Eng 9",\n     "credits": 1, "grade": "A"},\n    {"course": "Alg II/Trig",\n     "credits": 2, "grade": "A"},\n    ... ] }',
             font_size=12, color=ACCENT4)

# Warning card
warn_card = add_shape(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.2), BG_CARD, corner_radius=0.05)
add_shape(slide, Inches(0.8), Inches(5.8), Inches(0.1), Inches(1.2), ACCENT3)
icon_circle_text(slide, Inches(1.2), Inches(6.0), Inches(0.8), ACCENT3, "⚠️", font_size=24)
add_text_box(slide, Inches(2.2), Inches(6.05), Inches(9.8), Inches(0.8),
             "Design Note: Format heuristic, not content analysis.\nAny long message with newlines is treated as a transcript.",
             font_size=16, color=LIGHT_GRAY)


# ════════════════════════════════════════════
# SLIDE 4d — Short Text (Chat)
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8),
             "Input: Short Text (Chat)", font_size=36, color=WHITE, bold=True)
tag = add_shape(slide, Inches(6.5), Inches(0.45), Inches(3.0), Inches(0.55), ACCENT4, corner_radius=0.1)
tf = tag.text_frame
tf.paragraphs[0].text = "→ direct to pipeline"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.color.rgb = WHITE
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
add_shape(slide, Inches(0.8), Inches(1.15), Inches(2), Inches(0.05), ACCENT4)

# What qualifies - visual checklist
qual_card = add_shape(slide, Inches(0.8), Inches(1.7), Inches(5.2), Inches(2.8), BG_CARD, corner_radius=0.06)
add_text_box(slide, Inches(1.2), Inches(1.9), Inches(4.4), Inches(0.5),
             "✅  What Qualifies?", font_size=22, color=ACCENT4, bold=True)

quals = [
    ("❌", "No file attached"),
    ("❌", "NOT a URL (no http(s)://)"),
    ("❌", "≤ 300 chars OR no newlines"),
]
for i, (icon, text) in enumerate(quals):
    y = Inches(2.6) + i * Inches(0.7)
    icon_circle_text(slide, Inches(1.3), y, Inches(0.5), ACCENT3, icon, font_size=16)
    add_text_box(slide, Inches(2.0), y + Inches(0.07), Inches(3.6), Inches(0.45),
                 text, font_size=15, color=LIGHT_GRAY)

# What happens - visual
what_card = add_shape(slide, Inches(6.5), Inches(1.7), Inches(6.0), Inches(2.8), BG_CARD, corner_radius=0.06)
add_text_box(slide, Inches(6.9), Inches(1.9), Inches(5.2), Inches(0.5),
             "⚡  What Happens?", font_size=22, color=ACCENT, bold=True)

# Visual: message → pipeline
icon_circle_text(slide, Inches(7.2), Inches(2.7), Inches(1.0), ACCENT4, "💬", font_size=30)
add_text_box(slide, Inches(6.9), Inches(3.75), Inches(1.6), Inches(0.4),
             "Message", font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_arrow_right(slide, Inches(8.35), Inches(3.0), Inches(0.7), Inches(0.3), ACCENT4)

add_text_box(slide, Inches(9.2), Inches(2.55), Inches(3.0), Inches(0.5),
             "No extraction needed", font_size=14, color=ACCENT4, bold=True)
add_text_box(slide, Inches(9.2), Inches(3.0), Inches(3.0), Inches(0.5),
             "Passed directly to\nLangGraph pipeline", font_size=13, color=LIGHT_GRAY)

# ════════════════════════════════════════════
# SLIDE — Short Text Examples
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "Input: Short Text — Examples", font_size=36, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.15), Inches(2), Inches(0.05), ACCENT4)

ex_card = add_shape(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(5.3), BG_CARD, corner_radius=0.05)
add_text_box(slide, Inches(1.2), Inches(1.9), Inches(5), Inches(0.5),
             "💬  Example Questions → Pipeline Route", font_size=22, color=ACCENT4, bold=True)

# Column headers
add_text_box(slide, Inches(1.0), Inches(2.6), Inches(3.0), Inches(0.4),
             "Question", font_size=14, color=MID_GRAY, bold=True)
add_text_box(slide, Inches(4.5), Inches(2.6), Inches(2.0), Inches(0.4),
             "Route", font_size=14, color=MID_GRAY, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(7.0), Inches(2.6), Inches(4.0), Inches(0.4),
             "What Happens", font_size=14, color=MID_GRAY, bold=True)
add_shape(slide, Inches(1.0), Inches(3.0), Inches(11.0), Inches(0.02), MID_GRAY)

examples = [
    ("🧮", '"What is my CGPA?"', "calculation", "Runs CGPA calculator", ACCENT),
    ("🔀", '"Courses for AI track?"', "hybrid", "RAG + catalog lookup", ACCENT2),
    ("🔍", '"Explain the MAP program"', "retrieval", "Searches RAG KB", ACCENT3),
    ("👋", '"Hello!"', "reasoning", "General response", ACCENT4),
]
for i, (emoji, question, route, outcome, color) in enumerate(examples):
    x = Inches(1.0)
    y = Inches(3.3) + i * Inches(0.9)
    icon_circle_text(slide, x, y, Inches(0.6), color, emoji, font_size=20)
    add_text_box(slide, x + Inches(0.75), y + Inches(0.1), Inches(2.8), Inches(0.5),
                 question, font_size=16, color=WHITE)

    route_badge = add_shape(slide, Inches(4.5), y + Inches(0.08), Inches(1.8), Inches(0.45), color, corner_radius=0.08)
    tf = route_badge.text_frame
    tf.paragraphs[0].text = route
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_arrow_right(slide, Inches(6.45), y + Inches(0.12), Inches(0.45), Inches(0.25), MID_GRAY)

    add_text_box(slide, Inches(7.0), y + Inches(0.1), Inches(4.5), Inches(0.5),
                 outcome, font_size=16, color=LIGHT_GRAY)


# ════════════════════════════════════════════
# SLIDE — Stage 2: Section Title
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT2)

icon_circle_text(slide, Inches(5.65), Inches(1.2), Inches(2.0), ACCENT2, "⚙️", font_size=64)

add_text_box(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(0.5),
             "STAGE II", font_size=22, color=ACCENT2, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(4.0), Inches(10), Inches(1.0),
             "Data Processing", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(5.2), Inches(5.0), Inches(2.9), Inches(0.04), ACCENT2)
add_text_box(slide, Inches(1.5), Inches(5.3), Inches(10), Inches(0.6),
             "Linear LangGraph pipeline — all nodes run in sequence",
             font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

node_icons = [("📂", "Load Context", ACCENT), ("🏷️", "Classify Intent", ACCENT2),
              ("🔍", "Retrieve Docs", ACCENT3), ("🧮", "Calculate", ACCENT4)]
for i, (emoji, label, color) in enumerate(node_icons):
    x = Inches(1.8) + i * Inches(2.6)
    icon_circle_text(slide, x + Inches(0.3), Inches(6.0), Inches(0.8), color, emoji, font_size=24)
    add_text_box(slide, x - Inches(0.2), Inches(6.85), Inches(1.8), Inches(0.4),
                 label, font_size=13, color=color, bold=True, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════
# SLIDE — Node 1: load_context
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8),
             "Node: load_context", font_size=36, color=WHITE, bold=True)
tag = add_shape(slide, Inches(6.5), Inches(0.45), Inches(3.0), Inches(0.55), ACCENT, corner_radius=0.1)
tf = tag.text_frame
tf.paragraphs[0].text = "→ loads structured data"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.color.rgb = WHITE
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
add_shape(slide, Inches(0.8), Inches(1.15), Inches(2), Inches(0.05), ACCENT)

# What it does
what_card = add_shape(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(3.0), BG_CARD, corner_radius=0.06)
add_text_box(slide, Inches(1.2), Inches(1.9), Inches(4.8), Inches(0.5),
             "📂  What It Does", font_size=22, color=ACCENT, bold=True)
what_items = [
    ("📋", "Reads catalog_rules.json from disk"),
    ("📄", "Reads transcript_record.json from disk"),
    ("💾", "These were saved during Stage 1"),
    ("➡️", "Injects both into pipeline state"),
]
for i, (icon, text) in enumerate(what_items):
    y = Inches(2.6) + i * Inches(0.5)
    add_text_box(slide, Inches(1.3), y, Inches(0.5), Inches(0.4), icon, font_size=14)
    add_text_box(slide, Inches(1.9), y + Inches(0.02), Inches(4.0), Inches(0.4), text, font_size=15, color=LIGHT_GRAY)

# Two file cards on right
cat_card = add_shape(slide, Inches(6.8), Inches(1.7), Inches(5.5), Inches(1.3), BG_CARD, corner_radius=0.05)
add_shape(slide, Inches(6.8), Inches(1.7), Inches(5.5), Inches(0.06), ACCENT)
add_text_box(slide, Inches(7.1), Inches(1.85), Inches(4.8), Inches(0.4),
             "📋  catalog_rules.json", font_size=18, color=ACCENT, bold=True)
add_text_box(slide, Inches(7.1), Inches(2.3), Inches(4.8), Inches(0.5),
             "Program requirements, tracks, courses,\ngraduation rules — from URL scraping", font_size=14, color=LIGHT_GRAY)

trans_card = add_shape(slide, Inches(6.8), Inches(3.2), Inches(5.5), Inches(1.3), BG_CARD, corner_radius=0.05)
add_shape(slide, Inches(6.8), Inches(3.2), Inches(5.5), Inches(0.06), ACCENT2)
add_text_box(slide, Inches(7.1), Inches(3.35), Inches(4.8), Inches(0.4),
             "📄  transcript_record.json", font_size=18, color=ACCENT2, bold=True)
add_text_box(slide, Inches(7.1), Inches(3.8), Inches(4.8), Inches(0.5),
             "Student name, program, completed courses\nwith grades — from PDF/image/text upload", font_size=14, color=LIGHT_GRAY)

# Output card
out_card = add_shape(slide, Inches(0.8), Inches(5.1), Inches(11.5), Inches(2.1), BG_CARD, corner_radius=0.05)
add_shape(slide, Inches(0.8), Inches(5.1), Inches(11.5), Inches(0.06), ACCENT4)
add_text_box(slide, Inches(1.2), Inches(5.25), Inches(3), Inches(0.4),
             "⚡  Output → Pipeline State", font_size=18, color=ACCENT4, bold=True)
add_text_box(slide, Inches(1.2), Inches(5.75), Inches(10.8), Inches(1.2),
             '{ ...state,\n  "catalog_rules": { "program_name": "...", "required_credits": 120, "tracks": [...], ... },\n  "transcript_record": { "student_name": "...", "program": "...", "completed_courses": [...] } }',
             font_size=13, color=LIGHT_GRAY)


# ════════════════════════════════════════════
# SLIDE — Node 2: classify_intent
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8),
             "Node: classify_intent", font_size=36, color=WHITE, bold=True)
tag = add_shape(slide, Inches(6.5), Inches(0.45), Inches(3.5), Inches(0.55), ACCENT2, corner_radius=0.1)
tf = tag.text_frame
tf.paragraphs[0].text = "→ keyword-based routing"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.color.rgb = WHITE
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
add_shape(slide, Inches(0.8), Inches(1.15), Inches(2), Inches(0.05), ACCENT2)

# Two keyword lists
calc_card = add_shape(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(2.5), BG_CARD, corner_radius=0.06)
add_text_box(slide, Inches(1.2), Inches(1.9), Inches(4.8), Inches(0.5),
             "🧮  Calculation Keywords", font_size=20, color=ACCENT, bold=True)
add_text_box(slide, Inches(1.2), Inches(2.5), Inches(4.8), Inches(1.4),
             '"cgpa"  "gpa"  "credit"  "credits"\n"remaining"  "left"  "how many more"\n"graduate"  "eligible"',
             font_size=15, color=LIGHT_GRAY)

reason_card = add_shape(slide, Inches(6.8), Inches(1.7), Inches(5.5), Inches(2.5), BG_CARD, corner_radius=0.06)
add_text_box(slide, Inches(7.2), Inches(1.9), Inches(4.8), Inches(0.5),
             "💡  Reasoning Keywords", font_size=20, color=ACCENT2, bold=True)
add_text_box(slide, Inches(7.2), Inches(2.5), Inches(4.8), Inches(1.4),
             '"should i"  "which"  "recommend"\n"difference"  "compare"  "why"\n"advise"',
             font_size=15, color=LIGHT_GRAY)

# Routing logic
route_card = add_shape(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(2.8), BG_CARD, corner_radius=0.06)
add_text_box(slide, Inches(1.2), Inches(4.7), Inches(4), Inches(0.5),
             "🔀  Routing Logic", font_size=20, color=ACCENT3, bold=True)

routes = [
    ("Both match", "hybrid", "RAG + calculators run", ORANGE),
    ("Only calc match", "calculation", "Skip RAG, run calculators", ACCENT),
    ("Only reasoning match", "reasoning", "RAG search, skip calculators", ACCENT2),
    ("Neither match", "retrieval", "RAG search, skip calculators", ACCENT3),
]
for i, (condition, route, desc, color) in enumerate(routes):
    y = Inches(5.3) + i * Inches(0.47)
    add_text_box(slide, Inches(1.3), y, Inches(2.5), Inches(0.4),
                 condition, font_size=14, color=LIGHT_GRAY)
    route_badge = add_shape(slide, Inches(4.0), y, Inches(1.8), Inches(0.38), color, corner_radius=0.08)
    tf = route_badge.text_frame
    tf.paragraphs[0].text = route
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_arrow_right(slide, Inches(5.95), y + Inches(0.05), Inches(0.4), Inches(0.2), MID_GRAY)
    add_text_box(slide, Inches(6.5), y, Inches(5.0), Inches(0.4),
                 desc, font_size=14, color=LIGHT_GRAY)


# ════════════════════════════════════════════
# SLIDE — classify_intent Examples
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "classify_intent — Examples", font_size=36, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.15), Inches(2), Inches(0.05), ACCENT2)

ex_card = add_shape(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(5.5), BG_CARD, corner_radius=0.05)

# Column headers
add_text_box(slide, Inches(1.2), Inches(1.9), Inches(4.0), Inches(0.4),
             "Question", font_size=16, color=MID_GRAY, bold=True)
add_text_box(slide, Inches(5.5), Inches(1.9), Inches(1.5), Inches(0.4),
             "Calc?", font_size=16, color=MID_GRAY, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(7.0), Inches(1.9), Inches(1.5), Inches(0.4),
             "Reason?", font_size=16, color=MID_GRAY, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(9.0), Inches(1.9), Inches(2.5), Inches(0.4),
             "Route", font_size=16, color=MID_GRAY, bold=True, alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(1.2), Inches(2.35), Inches(10.8), Inches(0.02), MID_GRAY)

intent_examples = [
    ('"What is my CGPA?"', "YES", "no", "calculation", ACCENT),
    ('"How many credits remaining?"', "YES", "no", "calculation", ACCENT),
    ('"Should I take the AI track?"', "no", "YES", "reasoning", ACCENT2),
    ('"Which courses should I take?"', "no", "YES", "reasoning", ACCENT2),
    ('"Should I graduate now or take more credits?"', "YES", "YES", "hybrid", ORANGE),
    ('"Explain the MAP program"', "no", "no", "retrieval", ACCENT3),
    ('"Hello!"', "no", "no", "retrieval", ACCENT3),
]
for i, (question, calc, reason, route, color) in enumerate(intent_examples):
    y = Inches(2.6) + i * Inches(0.6)
    add_text_box(slide, Inches(1.2), y, Inches(4.0), Inches(0.5),
                 question, font_size=14, color=WHITE)
    add_text_box(slide, Inches(5.5), y, Inches(1.5), Inches(0.5),
                 calc, font_size=14, color=ACCENT4 if calc == "YES" else MID_GRAY,
                 bold=(calc == "YES"), alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(7.0), y, Inches(1.5), Inches(0.5),
                 reason, font_size=14, color=ACCENT4 if reason == "YES" else MID_GRAY,
                 bold=(reason == "YES"), alignment=PP_ALIGN.CENTER)
    route_badge = add_shape(slide, Inches(9.0), y + Inches(0.02), Inches(1.8), Inches(0.38), color, corner_radius=0.08)
    tf = route_badge.text_frame
    tf.paragraphs[0].text = route
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


# ════════════════════════════════════════════
# SLIDE — Node 3: retrieve_docs (RAG)
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8),
             "Node: retrieve_docs", font_size=36, color=WHITE, bold=True)
tag = add_shape(slide, Inches(6.0), Inches(0.45), Inches(4.0), Inches(0.55), ACCENT3, corner_radius=0.1)
tf = tag.text_frame
tf.paragraphs[0].text = "→ RAG semantic search"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.color.rgb = WHITE
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
add_shape(slide, Inches(0.8), Inches(1.15), Inches(2), Inches(0.05), ACCENT3)

# Skip condition
skip_card = add_shape(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.8), BG_CARD, corner_radius=0.05)
add_shape(slide, Inches(0.8), Inches(1.7), Inches(0.1), Inches(0.8), ORANGE)
icon_circle_text(slide, Inches(1.15), Inches(1.8), Inches(0.55), ORANGE, "⏭️", font_size=18)
add_text_box(slide, Inches(1.85), Inches(1.85), Inches(10.0), Inches(0.5),
             'Skipped entirely if route == "calculation" — returns empty chunks and citations',
             font_size=15, color=LIGHT_GRAY)

# Flow: Question → Embed → Cosine Search → Top 4 → Extract
flow_items = [
    ("💬", "Question", "User's original\nquestion text"),
    ("🔢", "Embed", "OpenAI embeddings\n(default model)"),
    ("📐", "Cosine\nSearch", "Compare against\nall doc chunks"),
    ("🏆", "Top 4\nResults", "similarity_top_k=4\nmost relevant"),
    ("✂️", "Extract", "Text[:1200] +\nsource filename"),
]
for i, (emoji, title, desc) in enumerate(flow_items):
    x = Inches(0.4) + i * Inches(2.55)
    card = add_shape(slide, x, Inches(2.8), Inches(2.2), Inches(2.8), BG_CARD, corner_radius=0.06)
    icon_circle_text(slide, x + Inches(0.6), Inches(2.95), Inches(1.0), ACCENT3, emoji, font_size=28)
    add_text_box(slide, x + Inches(0.1), Inches(4.05), Inches(2.0), Inches(0.6),
                 title, font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), Inches(4.55), Inches(2.0), Inches(0.8),
                 desc, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    if i < 4:
        add_arrow_right(slide, x + Inches(2.25), Inches(3.95), Inches(0.3), Inches(0.2), MID_GRAY)

# Output card
out_card = add_shape(slide, Inches(0.8), Inches(5.9), Inches(11.5), Inches(1.3), BG_CARD, corner_radius=0.05)
add_shape(slide, Inches(0.8), Inches(5.9), Inches(11.5), Inches(0.06), ACCENT4)
add_text_box(slide, Inches(1.2), Inches(6.0), Inches(3), Inches(0.4),
             "⚡  Output", font_size=18, color=ACCENT4, bold=True)
add_text_box(slide, Inches(1.2), Inches(6.4), Inches(10.8), Inches(0.6),
             '"retrieved_chunks": ["chunk text up to 1200 chars...", ...]    "citations": ["file_name_1.txt", "file_name_2.pdf"]',
             font_size=13, color=LIGHT_GRAY)


# ════════════════════════════════════════════
# SLIDE — RAG Knowledge Base Deep Dive
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "RAG Knowledge Base", font_size=36, color=WHITE, bold=True)
tag = add_shape(slide, Inches(6.0), Inches(0.45), Inches(3.5), Inches(0.55), ACCENT3, corner_radius=0.1)
tf = tag.text_frame
tf.paragraphs[0].text = "built once at startup"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.color.rgb = WHITE
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
add_shape(slide, Inches(0.8), Inches(1.15), Inches(2), Inches(0.05), ACCENT3)

# How the index is built — flow
build_card = add_shape(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(2.5), BG_CARD, corner_radius=0.06)
add_text_box(slide, Inches(1.2), Inches(1.9), Inches(5), Inches(0.5),
             "🏗️  How the Index Is Built (build_index)", font_size=20, color=ACCENT2, bold=True)

build_steps = [
    ("📁", "Scan\nFiles", "processed_data/\n.pdf .txt .md"),
    ("📖", "Read &\nParse", "SimpleDirectoryReader\n+ PDFReader"),
    ("✂️", "Chunk", "Split docs into\nsmaller pieces"),
    ("🔢", "Embed", "OpenAI embeddings\nfor each chunk"),
    ("🗄️", "Store", "In-memory\nvector index"),
]
for i, (emoji, title, desc) in enumerate(build_steps):
    x = Inches(1.0) + i * Inches(2.2)
    icon_circle_text(slide, x + Inches(0.35), Inches(2.5), Inches(0.75), ACCENT3, emoji, font_size=22)
    add_text_box(slide, x, Inches(3.3), Inches(1.5), Inches(0.45),
                 title, font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(3.7), Inches(1.5), Inches(0.4),
                 desc, font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    if i < 4:
        add_arrow_right(slide, x + Inches(1.55), Inches(2.7), Inches(0.3), Inches(0.18), MID_GRAY)

# What's in the knowledge base
content_card = add_shape(slide, Inches(0.8), Inches(4.5), Inches(5.5), Inches(2.8), BG_CARD, corner_radius=0.06)
add_text_box(slide, Inches(1.2), Inches(4.7), Inches(4.8), Inches(0.5),
             "📚  Contents (processed_data/)", font_size=18, color=ACCENT, bold=True)
content_items = [
    ("📄", "Course lecture slides (CS 4440/5440)"),
    ("📝", "Prompt engineering papers"),
    ("🎬", "Video lecture transcripts"),
    ("📘", "AI/ML reference documents"),
]
for i, (icon, text) in enumerate(content_items):
    y = Inches(5.3) + i * Inches(0.45)
    add_text_box(slide, Inches(1.3), y, Inches(0.5), Inches(0.4), icon, font_size=14)
    add_text_box(slide, Inches(1.9), y + Inches(0.02), Inches(4.0), Inches(0.4), text, font_size=14, color=LIGHT_GRAY)

# Key properties
props_card = add_shape(slide, Inches(6.8), Inches(4.5), Inches(5.5), Inches(2.8), BG_CARD, corner_radius=0.06)
add_text_box(slide, Inches(7.2), Inches(4.7), Inches(4.8), Inches(0.5),
             "🔑  Key Properties", font_size=18, color=ACCENT4, bold=True)
prop_items = [
    ("🔒", "Static — never modified by uploads"),
    ("💨", "In-memory — rebuilt on each restart"),
    ("🔍", "similarity_top_k=4 per query"),
    ("✂️", "Chunks capped at 1200 chars each"),
]
for i, (icon, text) in enumerate(prop_items):
    y = Inches(5.3) + i * Inches(0.45)
    add_text_box(slide, Inches(7.3), y, Inches(0.5), Inches(0.4), icon, font_size=14)
    add_text_box(slide, Inches(7.9), y + Inches(0.02), Inches(4.0), Inches(0.4), text, font_size=14, color=LIGHT_GRAY)


# ════════════════════════════════════════════
# SLIDE — Node 4: calculate_tools
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8),
             "Node: calculate_tools", font_size=36, color=WHITE, bold=True)
tag = add_shape(slide, Inches(6.0), Inches(0.45), Inches(3.5), Inches(0.55), ACCENT4, corner_radius=0.1)
tf = tag.text_frame
tf.paragraphs[0].text = "→ phrase-matched calcs"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.color.rgb = WHITE
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
add_shape(slide, Inches(0.8), Inches(1.15), Inches(2), Inches(0.05), ACCENT4)

# Three calculator cards
calc_cards = [
    ("🧮", "CGPA Calculator", '"my cgpa" / "my gpa"',
     "For each course:\n  quality_pts = credits × grade_pts\nCGPA = total_quality / total_credits",
     ACCENT),
    ("📊", "Remaining Credits", '"remaining credits" / "credits left"\n/ "how many more credits"',
     "remaining = required_credits\n  - sum(completed credits)\nUses catalog_rules.required_credits",
     ACCENT2),
    ("🎓", "Graduation Status", '"can i graduate" /\n"am i eligible to graduate"',
     "Checks: remaining_credits <= 0\n  AND cgpa >= minimum_cgpa\nReturns eligible: true/false",
     ACCENT3),
]
for i, (emoji, title, trigger, logic, color) in enumerate(calc_cards):
    x = Inches(0.5) + i * Inches(4.2)
    card = add_shape(slide, x, Inches(1.7), Inches(3.8), Inches(5.3), BG_CARD, corner_radius=0.06)
    add_shape(slide, x, Inches(1.7), Inches(3.8), Inches(0.07), color)
    icon_circle_text(slide, x + Inches(1.3), Inches(1.95), Inches(1.2), color, emoji, font_size=34)
    add_text_box(slide, x + Inches(0.15), Inches(3.3), Inches(3.5), Inches(0.5),
                 title, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_shape(slide, x + Inches(1.0), Inches(3.8), Inches(1.8), Inches(0.03), color)

    trig_card = add_shape(slide, x + Inches(0.15), Inches(4.0), Inches(3.5), Inches(0.9), BG_CARD_LIGHT, corner_radius=0.04)
    add_text_box(slide, x + Inches(0.25), Inches(4.05), Inches(1.2), Inches(0.3),
                 "TRIGGER", font_size=10, color=color, bold=True)
    add_text_box(slide, x + Inches(0.25), Inches(4.3), Inches(3.3), Inches(0.5),
                 trigger, font_size=11, color=LIGHT_GRAY)

    add_text_box(slide, x + Inches(0.15), Inches(5.1), Inches(3.5), Inches(1.7),
                 logic, font_size=12, color=LIGHT_GRAY)


# ════════════════════════════════════════════
# SLIDE — calculate_tools Output Examples
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "calculate_tools — Output Examples", font_size=36, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.15), Inches(2), Inches(0.05), ACCENT4)

# CGPA output
cgpa_card = add_shape(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(1.7), BG_CARD, corner_radius=0.05)
add_shape(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.06), ACCENT)
add_text_box(slide, Inches(1.2), Inches(1.85), Inches(3), Inches(0.4),
             "🧮  CGPA Result", font_size=18, color=ACCENT, bold=True)
add_text_box(slide, Inches(1.2), Inches(2.3), Inches(10.8), Inches(0.9),
             '{ "type": "cgpa", "total_credits": 45.0, "total_quality_points": 162.0, "cgpa": 3.6,\n  "breakdown": [{"course": "Chemistry", "credits": 1, "grade": "A", "grade_points": 4.0, "quality_points": 4.0}, ...] }',
             font_size=12, color=LIGHT_GRAY)

# Remaining credits output
rem_card = add_shape(slide, Inches(0.8), Inches(3.7), Inches(11.5), Inches(1.5), BG_CARD, corner_radius=0.05)
add_shape(slide, Inches(0.8), Inches(3.7), Inches(11.5), Inches(0.06), ACCENT2)
add_text_box(slide, Inches(1.2), Inches(3.85), Inches(3), Inches(0.4),
             "📊  Remaining Credits Result", font_size=18, color=ACCENT2, bold=True)
add_text_box(slide, Inches(1.2), Inches(4.3), Inches(10.8), Inches(0.7),
             '{ "type": "remaining_credits", "completed_credits": 45.0, "required_credits": 120.0, "remaining_credits": 75.0 }',
             font_size=13, color=LIGHT_GRAY)

# Graduation status output
grad_card = add_shape(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.7), BG_CARD, corner_radius=0.05)
add_shape(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.06), ACCENT3)
add_text_box(slide, Inches(1.2), Inches(5.65), Inches(3), Inches(0.4),
             "🎓  Graduation Status Result", font_size=18, color=ACCENT3, bold=True)
add_text_box(slide, Inches(1.2), Inches(6.1), Inches(10.8), Inches(0.9),
             '{ "type": "graduation_status", "cgpa": 3.6, "minimum_cgpa_for_graduation": 2.0,\n  "completed_credits": 45.0, "required_credits": 120.0, "remaining_credits": 75.0, "eligible": false }',
             font_size=13, color=LIGHT_GRAY)


# ════════════════════════════════════════════
# SLIDE — Two Data Stores
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "Two Separate Data Stores", font_size=38, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.05), ACCENT)

# Vector DB visual
v_card = add_shape(slide, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.2), BG_CARD, corner_radius=0.06)
add_shape(slide, Inches(0.8), Inches(1.7), Inches(5.6), Inches(0.07), ACCENT2)
icon_circle_text(slide, Inches(2.9), Inches(1.95), Inches(1.2), ACCENT2, "🗄️", font_size=38)
add_text_box(slide, Inches(1.2), Inches(3.3), Inches(4.8), Inches(0.5),
             "Vector Index (RAG KB)", font_size=22, color=ACCENT2, bold=True, alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(2.0), Inches(3.8), Inches(2.8), Inches(0.03), ACCENT2)

v_items = [
    ("📁", "Built from processed_data/ at startup"),
    ("📚", "Course lectures, papers, transcripts"),
    ("🔢", "Embedded via OpenAI embeddings"),
    ("🔍", "Queried by semantic similarity (top 4)"),
    ("🔒", "Static — uploads don't modify it"),
    ("💨", "In-memory — rebuilt on each restart"),
]
for i, (icon, text) in enumerate(v_items):
    y = Inches(4.05) + i * Inches(0.5)
    add_text_box(slide, Inches(1.3), y, Inches(0.5), Inches(0.4), icon, font_size=14)
    add_text_box(slide, Inches(1.8), y + Inches(0.02), Inches(4.2), Inches(0.4), text, font_size=13, color=LIGHT_GRAY)

# Structured JSON visual
j_card = add_shape(slide, Inches(6.9), Inches(1.7), Inches(5.6), Inches(5.2), BG_CARD, corner_radius=0.06)
add_shape(slide, Inches(6.9), Inches(1.7), Inches(5.6), Inches(0.07), ACCENT)
icon_circle_text(slide, Inches(9.0), Inches(1.95), Inches(1.2), ACCENT, "📋", font_size=38)
add_text_box(slide, Inches(7.3), Inches(3.3), Inches(4.8), Inches(0.5),
             "Structured JSON", font_size=22, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(8.1), Inches(3.8), Inches(2.8), Inches(0.03), ACCENT)

j_items = [
    ("🔗", "catalog_rules.json — from URLs"),
    ("📄", "transcript_record.json — from uploads"),
    ("🤖", "Extracted at runtime by GPT-4.1"),
    ("♻️", "Overwritten each time (single-user)"),
    ("🧮", "Used by: calculators + answer prompt"),
    ("💾", "Persisted in extracted_data/ dir"),
]
for i, (icon, text) in enumerate(j_items):
    y = Inches(4.05) + i * Inches(0.5)
    add_text_box(slide, Inches(7.4), y, Inches(0.5), Inches(0.4), icon, font_size=14)
    add_text_box(slide, Inches(7.9), y + Inches(0.02), Inches(4.2), Inches(0.4), text, font_size=13, color=LIGHT_GRAY)


# ════════════════════════════════════════════
# SLIDE 7 — Stage 3: Answer Generation
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "Stage 3: Answer Generation", font_size=38, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.05), ACCENT4)

# Four source cards feeding into GPT-4.1
sources = [
    ("🔍", "RAG Chunks", "Semantic search\nresults", ACCENT2),
    ("📋", "Catalog JSON", "Program rules\n& requirements", ACCENT),
    ("📄", "Transcript", "Student record\n& grades", ACCENT3),
    ("🧮", "Calc Results", "CGPA, credits,\ngrad status", ACCENT4),
]

for i, (emoji, title, desc, color) in enumerate(sources):
    x = Inches(0.6) + i * Inches(3.15)
    card = add_shape(slide, x, Inches(1.6), Inches(2.8), Inches(2.3), BG_CARD, corner_radius=0.06)
    add_shape(slide, x, Inches(1.6), Inches(2.8), Inches(0.06), color)
    icon_circle_text(slide, x + Inches(0.85), Inches(1.75), Inches(1.1), color, emoji, font_size=30)
    add_text_box(slide, x + Inches(0.1), Inches(2.95), Inches(2.6), Inches(0.4),
                 title, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), Inches(3.35), Inches(2.6), Inches(0.5),
                 desc, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Arrows down
for i in range(4):
    x = Inches(1.4) + i * Inches(3.15)
    add_arrow_down(slide, x, Inches(3.95), Inches(0.3), Inches(0.5), MID_GRAY)

# Big GPT-4.1 box
gpt_card = add_shape(slide, Inches(1.5), Inches(4.6), Inches(10.3), Inches(1.4), ACCENT, corner_radius=0.08)
icon_circle_text(slide, Inches(2.0), Inches(4.75), Inches(1.1), DARK_ACCENT, "🤖", font_size=34)
add_text_box(slide, Inches(3.3), Inches(4.75), Inches(3), Inches(0.5),
             "GPT-4.1", font_size=28, color=WHITE, bold=True)
add_text_box(slide, Inches(3.3), Inches(5.3), Inches(8), Inches(0.5),
             "Assembles all context into one prompt → generates grounded, accurate response",
             font_size=15, color=RGBColor(0xDD, 0xDD, 0xFF))

# Arrow down to response
add_arrow_down(slide, Inches(6.4), Inches(6.05), Inches(0.3), Inches(0.5), ACCENT4)

# Response
resp_card = add_shape(slide, Inches(4.0), Inches(6.6), Inches(5.3), Inches(0.7), ACCENT4, corner_radius=0.08)
tf = resp_card.text_frame
tf.paragraphs[0].text = "✨  Student gets a grounded, accurate answer"
tf.paragraphs[0].font.size = Pt(16)
tf.paragraphs[0].font.color.rgb = WHITE
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.vertical_anchor = MSO_ANCHOR.MIDDLE


# ════════════════════════════════════════════
# SLIDE 8 — Tech Stack
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "Tech Stack", font_size=38, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.05), ACCENT)

stack_items = [
    ("🔀", "Orchestration", "LangGraph", "Controls multi-stage pipeline", ACCENT),
    ("🤖", "LLM", "GPT-4.1 (OpenAI)", "Extraction & answer generation", ACCENT2),
    ("🔍", "RAG", "Vector Index + Search", "Retrieves knowledge chunks", ACCENT3),
    ("⚡", "Backend", "FastAPI", "REST API for uploads & chat", ACCENT4),
    ("🖥️", "Frontend", "Web Interface", "Chat UI for student interaction", ORANGE),
]

for i, (emoji, category, tech, desc, color) in enumerate(stack_items):
    y = Inches(1.6) + i * Inches(1.15)
    card = add_shape(slide, Inches(0.8), y, Inches(11.5), Inches(0.95), BG_CARD, corner_radius=0.05)
    add_shape(slide, Inches(0.8), y, Inches(0.1), Inches(0.95), color)

    icon_circle_text(slide, Inches(1.2), y + Inches(0.12), Inches(0.7), color, emoji, font_size=22)

    add_text_box(slide, Inches(2.2), y + Inches(0.2), Inches(2.0), Inches(0.55),
                 category, font_size=15, color=MID_GRAY)
    add_text_box(slide, Inches(4.2), y + Inches(0.15), Inches(3.5), Inches(0.6),
                 tech, font_size=20, color=WHITE, bold=True)
    add_text_box(slide, Inches(8.0), y + Inches(0.2), Inches(4.0), Inches(0.55),
                 desc, font_size=15, color=LIGHT_GRAY)


# ════════════════════════════════════════════
# SLIDE 9 — Key Design Characteristics
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "Key Design Characteristics", font_size=38, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.05), ACCENT)

characteristics = [
    ("🔗", "LangGraph Controlled", "Entire system orchestrated\nby a LangGraph StateGraph", ACCENT),
    ("📂", "Format Routing", "Input type determines path\nNo content analysis", ACCENT2),
    ("🏷️", "Keyword Intent", "Substring matching\nNot LLM-based", ACCENT3),
    ("🔒", "Static RAG KB", "Built once at startup\nUploads don't extend it", ACCENT4),
    ("👤", "Single-User", "One transcript + catalog\nEach upload overwrites", ORANGE),
    ("🎯", "Grounded Answers", "RAG + structured data\n+ calcs = less hallucination", ACCENT),
]

for i, (emoji, title, desc, color) in enumerate(characteristics):
    x = Inches(0.6) + (i % 3) * Inches(4.2)
    y = Inches(1.6) + (i // 3) * Inches(2.85)
    card = add_shape(slide, x, y, Inches(3.8), Inches(2.5), BG_CARD, corner_radius=0.06)
    icon_circle_text(slide, x + Inches(1.35), y + Inches(0.2), Inches(1.1), color, emoji, font_size=32)
    add_text_box(slide, x + Inches(0.15), y + Inches(1.4), Inches(3.5), Inches(0.45),
                 title, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.15), y + Inches(1.85), Inches(3.5), Inches(0.6),
                 desc, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════
# SLIDE — Interface: Welcome Screen
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "Project Interface", font_size=38, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.05), ACCENT)

# Browser window mockup
browser_bg = add_shape(slide, Inches(1.5), Inches(1.6), Inches(10.3), Inches(5.5), RGBColor(0x0A, 0x0A, 0x0A), corner_radius=0.04)
browser_bg.line.color.rgb = MID_GRAY
browser_bg.line.width = Pt(1)

# Top bar
add_shape(slide, Inches(1.5), Inches(1.6), Inches(10.3), Inches(0.65), RGBColor(0x17, 0x17, 0x17))
# Window dots
dot_colors = [RGBColor(0xFF, 0x5F, 0x57), RGBColor(0xFF, 0xBD, 0x2E), RGBColor(0x28, 0xCA, 0x41)]
for j, dc in enumerate(dot_colors):
    d = add_circle(slide, Inches(1.75) + j * Inches(0.3), Inches(1.73), Inches(0.18), dc)
add_text_box(slide, Inches(4.5), Inches(1.68), Inches(4), Inches(0.5),
             "Smart Chat Bot", font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Clear chat button mockup
clear_btn = add_shape(slide, Inches(10.2), Inches(1.72), Inches(1.2), Inches(0.38), BG_CARD, corner_radius=0.1)
clear_btn.line.color.rgb = MID_GRAY
clear_btn.line.width = Pt(1)
tf = clear_btn.text_frame
tf.paragraphs[0].text = "Clear chat"
tf.paragraphs[0].font.size = Pt(10)
tf.paragraphs[0].font.color.rgb = MID_GRAY
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

# Welcome content area
add_text_box(slide, Inches(3.5), Inches(3.0), Inches(6.0), Inches(0.7),
             "How can I help with your academics?",
             font_size=24, color=RGBColor(0x00, 0x00, 0x00), bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(3.0), Inches(3.8), Inches(7.0), Inches(1.2),
             'Type a question, paste a program catalog URL, or upload\na transcript file with the + button.',
             font_size=13, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Composer area mockup
composer_bg = add_shape(slide, Inches(3.0), Inches(5.6), Inches(7.0), Inches(0.65), RGBColor(0x2F, 0x2F, 0x2F), corner_radius=0.15)
composer_bg.line.color.rgb = RGBColor(0x4A, 0x4A, 0x52)
composer_bg.line.width = Pt(1)

# + button
plus_btn = add_circle(slide, Inches(3.15), Inches(5.68), Inches(0.5), BG_CARD)
tf = plus_btn.text_frame
tf.paragraphs[0].text = "+"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.color.rgb = MID_GRAY
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

add_text_box(slide, Inches(3.8), Inches(5.73), Inches(5.0), Inches(0.4),
             "Ask anything, paste a URL, or upload a file...",
             font_size=12, color=MID_GRAY)

# Send button
send_btn = add_circle(slide, Inches(9.35), Inches(5.68), Inches(0.5), RGBColor(0x10, 0xA3, 0x7F))
tf = send_btn.text_frame
tf.paragraphs[0].text = "↑"
tf.paragraphs[0].font.size = Pt(16)
tf.paragraphs[0].font.color.rgb = WHITE
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

# Callout labels
callouts = [
    (Inches(0.2), Inches(2.0), "📎 Upload files\n(PDF, images)", ACCENT2),
    (Inches(0.2), Inches(3.6), "🔗 Paste catalog URLs", ACCENT),
    (Inches(0.2), Inches(5.2), "💬 Ask questions\nin natural language", ACCENT4),
]
for cx, cy, text, color in callouts:
    add_text_box(slide, cx, cy, Inches(1.6), Inches(0.7), text, font_size=11, color=color, bold=True)


# ════════════════════════════════════════════
# SLIDE — Interface: Example Conversation
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "Interface — Example Conversation", font_size=38, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.05), ACCENT)

# Chat window mockup
chat_bg = add_shape(slide, Inches(1.5), Inches(1.6), Inches(7.5), Inches(5.5), RGBColor(0x0A, 0x0A, 0x0A), corner_radius=0.04)
chat_bg.line.color.rgb = MID_GRAY
chat_bg.line.width = Pt(1)

# Top bar
add_shape(slide, Inches(1.5), Inches(1.6), Inches(7.5), Inches(0.5), RGBColor(0x17, 0x17, 0x17))
add_text_box(slide, Inches(3.5), Inches(1.62), Inches(3.5), Inches(0.4),
             "Smart Chat Bot", font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Message 1: User sends URL
add_text_box(slide, Inches(1.8), Inches(2.3), Inches(1.0), Inches(0.3),
             "You", font_size=9, color=MID_GRAY)
msg1 = add_shape(slide, Inches(1.8), Inches(2.55), Inches(6.5), Inches(0.55), RGBColor(0x30, 0x30, 0x30), corner_radius=0.08)
tf = msg1.text_frame
tf.paragraphs[0].text = "https://catalogs.nmsu.edu/nmsu/arts-sciences/computer-science/"
tf.paragraphs[0].font.size = Pt(10)
tf.paragraphs[0].font.color.rgb = WHITE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

# Message 2: Assistant confirms extraction
add_text_box(slide, Inches(1.8), Inches(3.25), Inches(1.2), Inches(0.3),
             "Assistant", font_size=9, color=MID_GRAY)
msg2 = add_shape(slide, Inches(1.8), Inches(3.5), Inches(6.5), Inches(0.7), BG_CARD, corner_radius=0.08)
msg2.line.color.rgb = RGBColor(0x2A, 0x2A, 0x2F)
msg2.line.width = Pt(1)
tf = msg2.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "I've extracted the catalog rules for Computer Science. Found 120 required credits, AI/Cybersecurity tracks, and 45 required courses."
tf.paragraphs[0].font.size = Pt(9)
tf.paragraphs[0].font.color.rgb = LIGHT_GRAY
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

# Message 3: User asks CGPA
add_text_box(slide, Inches(1.8), Inches(4.35), Inches(1.0), Inches(0.3),
             "You", font_size=9, color=MID_GRAY)
msg3 = add_shape(slide, Inches(1.8), Inches(4.6), Inches(6.5), Inches(0.45), RGBColor(0x30, 0x30, 0x30), corner_radius=0.08)
tf = msg3.text_frame
tf.paragraphs[0].text = "What is my CGPA?"
tf.paragraphs[0].font.size = Pt(10)
tf.paragraphs[0].font.color.rgb = WHITE
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

# Message 4: Assistant calculates
add_text_box(slide, Inches(1.8), Inches(5.2), Inches(1.2), Inches(0.3),
             "Assistant", font_size=9, color=MID_GRAY)
msg4 = add_shape(slide, Inches(1.8), Inches(5.45), Inches(6.5), Inches(0.7), BG_CARD, corner_radius=0.08)
msg4.line.color.rgb = RGBColor(0x2A, 0x2A, 0x2F)
msg4.line.width = Pt(1)
tf = msg4.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "Based on your transcript, your cumulative GPA is 3.60 across 45 completed credits. You have 75 credits remaining to graduate."
tf.paragraphs[0].font.size = Pt(9)
tf.paragraphs[0].font.color.rgb = LIGHT_GRAY
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

# Annotations on the right side showing what happens behind the scenes
ann_y_positions = [
    (Inches(2.55), "→ URL detected → scrape → GPT-4.1\n   extracts catalog_rules.json", ACCENT),
    (Inches(3.5), "→ LangGraph: load_context\n   + generate_answer", ACCENT2),
    (Inches(4.6), '→ LangGraph: classify_intent\n   routes to "calculation"', ACCENT3),
    (Inches(5.45), "→ LangGraph: calculate_tools\n   runs CGPA calc → generate_answer", ACCENT4),
]
add_text_box(slide, Inches(9.3), Inches(1.8), Inches(3.5), Inches(0.5),
             "⚡ Behind the Scenes", font_size=16, color=ACCENT, bold=True)
add_shape(slide, Inches(9.3), Inches(2.25), Inches(2.0), Inches(0.03), ACCENT)
for ay, text, color in ann_y_positions:
    add_text_box(slide, Inches(9.3), ay, Inches(3.8), Inches(0.8),
                 text, font_size=10, color=color)


# ════════════════════════════════════════════
# SLIDE — Thank You
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_shape(slide, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), ACCENT)

# Decorative circles
add_circle(slide, Inches(12.3), Inches(-0.5), Inches(1.5), DARK_ACCENT)
add_circle(slide, Inches(-0.3), Inches(6.5), Inches(1.5), DARK_ACCENT)

icon_circle_text(slide, Inches(5.65), Inches(1.3), Inches(2.0), ACCENT, "🎓", font_size=64)

add_text_box(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(1.0),
             "Thank You", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(5.2), Inches(4.5), Inches(2.9), Inches(0.04), ACCENT)

add_text_box(slide, Inches(1.5), Inches(4.8), Inches(10), Inches(0.7),
             "Smart Academic Assistant", font_size=24, color=ACCENT2, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(5.6), Inches(10), Inches(0.6),
             "Questions?", font_size=22, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


# ── Save ──
output_path = "/home/grad24/mnguyen/smart-genai-bot/Smart_Academic_Assistant_6.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
