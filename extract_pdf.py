"""
Document extractor for RAG ingestion. Handles PDF, PPTX, and TXT.

Usage:
    python extract_pdf.py path/to/file.pdf
    python extract_pdf.py path/to/slides.pptx
    python extract_pdf.py path/to/notes.txt

Output is always saved to processed_data/<stem>.txt so ingest.py can pick it up.
"""

import argparse
import os
import sys
from pathlib import Path

from dotenv import load_dotenv
from openai import OpenAI
from pypdf import PdfReader

load_dotenv()

client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

PROJECT_ROOT = Path(__file__).parent
OUT_DIR = PROJECT_ROOT / "processed_data"

SYSTEM_PROMPT = """You are a precise document extraction assistant. Your job is to extract and structure content from a single page or slide for use in a Retrieval-Augmented Generation (RAG) system.

Rules:
- Extract ALL content — leave nothing out. More is better.
- Use labeled semantic sections so a retriever can identify what type of content each block contains.
- Common section labels (use whichever apply, add others as needed):
    TITLE, AUTHORS, AFFILIATIONS, ABSTRACT, KEYWORDS, SECTION X - NAME,
    SUBSECTION X.Y - NAME, DEFINITION, CONCEPT, TECHNIQUE, DESCRIPTION,
    PIPELINE, ARCHITECTURE, ALGORITHM, FORMULA, TABLE, FIGURE DESCRIPTION,
    EXAMPLE INPUT, EXAMPLE OUTPUT, KEY INSIGHTS, PERFORMANCE, LIMITATIONS,
    REFERENCES, FOOTNOTES
- Preserve tables verbatim using ASCII alignment.
- Preserve formulas and equations exactly as they appear.
- Preserve example prompts/outputs verbatim — wrap them in EXAMPLE INPUT / EXAMPLE OUTPUT.
- Do NOT summarize, paraphrase, or omit any content.
- Do NOT add commentary or opinions.
- Output plain text only — no markdown formatting, no bullet symbols.
- If content continues from the previous page, label it with (CONTINUED).
"""

UNIT_PROMPT = """Extract all content from this document unit.

Filename: {filename}
Unit {page_num} of {total_pages}

Raw text:
---
{raw_text}
---

Structure the content with labeled semantic sections. Start directly with the content — no preamble."""


# ── Shared helpers ────────────────────────────────────────────────────────────

def _structure_unit(filename: str, page_num: int, total: int, raw_text: str) -> str:
    if not raw_text.strip():
        return "(No extractable text — may be an image or blank.)"
    response = client.chat.completions.create(
        model="gpt-4.1",
        messages=[
            {"role": "system", "content": SYSTEM_PROMPT},
            {"role": "user", "content": UNIT_PROMPT.format(
                filename=filename, page_num=page_num, total_pages=total, raw_text=raw_text
            )},
        ],
        temperature=0.1,
    )
    return response.choices[0].message.content.strip()


def _build_output(filename: str, raw_units: list) -> str:
    total = len(raw_units)
    sep = "=" * 80
    lines = []
    for i, raw_text in enumerate(raw_units, start=1):
        print(f"  {i}/{total}...", end="", flush=True)
        structured = _structure_unit(filename, i, total, raw_text)
        lines += [sep, f"PAGE {i}", sep, f"SOURCE: {filename} | Page {i}", "", structured, ""]
        print(" done")
    return "\n".join(lines)


# ── Per-format handlers ───────────────────────────────────────────────────────

def _handle_pdf(path: Path) -> str:
    reader = PdfReader(str(path))
    units = [(page.extract_text() or "").strip() for page in reader.pages]
    return _build_output(path.name, units)


def _collect_shape_text(shape) -> list:
    texts = []
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            t = para.text.strip()
            if t:
                texts.append(t)
    if shape.has_table:
        for row in shape.table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                texts.append(" | ".join(cells))
    if shape.shape_type == 6:  # GROUP
        for child in shape.shapes:
            texts.extend(_collect_shape_text(child))
    return texts


def _handle_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    units = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            texts.extend(_collect_shape_text(shape))
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                texts.append(f"SPEAKER NOTES: {notes}")
        units.append("\n".join(texts))
    return _build_output(path.name, units)


def _handle_txt(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


# ── Dispatch table ────────────────────────────────────────────────────────────

HANDLERS = {
    ".pdf": _handle_pdf,
    ".pptx": _handle_pptx,
    ".txt": _handle_txt,
}


def extract_document(path_str: str) -> str:
    path = Path(path_str).resolve()
    if not path.exists():
        print(f"Error: file not found: {path}", file=sys.stderr)
        sys.exit(1)

    ext = path.suffix.lower()
    handler = HANDLERS.get(ext)
    if not handler:
        print(f"Error: unsupported file type: {ext}", file=sys.stderr)
        sys.exit(1)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    output_path = OUT_DIR / f"{path.stem}.txt"

    print(f"Extracting: {path.name}")
    print(f"Output: {output_path}")

    content = handler(path)
    output_path.write_text(content, encoding="utf-8")
    print(f"\nSaved -> {output_path}")
    return str(output_path)


def main():
    parser = argparse.ArgumentParser(
        description="Extract document content into processed_data/ for RAG ingestion"
    )
    parser.add_argument("file", help="Path to file (.pdf, .pptx, or .txt)")
    args = parser.parse_args()
    extract_document(args.file)


if __name__ == "__main__":
    main()
